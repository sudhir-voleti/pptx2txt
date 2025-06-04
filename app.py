import streamlit as st
from pptx import Presentation
import tempfile
import os

# Streamlit app title
st.title("PPTX Text Extractor")
st.write("Upload a PowerPoint (.pptx) file to extract its text content into a downloadable .txt file.")

# File uploader for .pptx files
uploaded_file = st.file_uploader("Choose a .pptx file", type="pptx")

if uploaded_file is not None:
    # Save the uploaded file to a temporary location
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
        tmp_file.write(uploaded_file.read())
        tmp_pptx_path = tmp_file.name

    try:
        # Load the presentation
        prs = Presentation(tmp_pptx_path)
        
        # Extract text from all slides
        extracted_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):  # Check if the shape has text
                    extracted_text.append(shape.text.strip())
        
        # Join the extracted text with newlines
        final_text = "\n".join(text for text in extracted_text if text)

        # Save the extracted text to a temporary .txt file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".txt") as tmp_txt:
            tmp_txt.write(final_text.encode('utf-8'))
            tmp_txt_path = tmp_txt.name

        # Provide a download button for the .txt file
        with open(tmp_txt_path, "rb") as file:
            st.download_button(
                label="Download Extracted Text (.txt)",
                data=file,
                file_name="extracted_pptx_text.txt",
                mime="text/plain"
            )

        # Display the extracted text in the app for preview
        st.subheader("Extracted Text Preview")
        st.text_area("Preview", final_text, height=200)

    except Exception as e:
        st.error(f"An error occurred: {str(e)}")

    finally:
        # Clean up temporary files
        os.unlink(tmp_pptx_path)
        if 'tmp_txt_path' in locals():
            os.unlink(tmp_txt_path)

st.write("Note: This app extracts text from slides, ignoring images, shapes without text, and other non-text elements.")
